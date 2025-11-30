import glob
import json
from pathlib import Path
from typing import List, Dict, Tuple

import numpy as np
import faiss
from sentence_transformers import SentenceTransformer
from pypdf import PdfReader
from pptx import Presentation

from .config import DATASET_DIR, INDEX_DIR, EMBED_MODEL

INDEX_PATH = INDEX_DIR / "faiss.index"
CHUNKS_PATH = INDEX_DIR / "chunks.jsonl"
META_PATH = INDEX_DIR / "meta.json"

def _chunk_text(text: str, size: int = 900, overlap: int = 180) -> List[str]:
    text = (text or "").replace("\x00", "")
    chunks = []
    i = 0
    while i < len(text):
        chunk = text[i : i + size]
        chunks.append(chunk)
        i += max(1, size - overlap)
    return [c.strip() for c in chunks if c.strip()]

def _read_pdf(path: Path) -> List[Tuple[str, Dict]]:
    items = []
    reader = PdfReader(str(path))
    for i, page in enumerate(reader.pages, start=1):
        txt = page.extract_text() or ""
        for c in _chunk_text(txt):
            items.append((c, {"source": str(path), "page": i, "doc_type": ".pdf"}))
    return items

def _read_pptx(path: Path) -> List[Tuple[str, Dict]]:
    items = []
    prs = Presentation(str(path))
    for i, slide in enumerate(prs.slides, start=1):
        buf = []
        for shp in slide.shapes:
            if hasattr(shp, "text") and shp.text:
                buf.append(shp.text)
        txt = "\n".join(buf)
        for c in _chunk_text(txt):
            items.append((c, {"source": str(path), "page": i, "doc_type": ".pptx"}))
    return items

def build_index() -> Dict:
    texts, metas = [], []
    pdfs = glob.glob(str(DATASET_DIR / "**" / "*.pdf"), recursive=True)
    pptxs = glob.glob(str(DATASET_DIR / "**" / "*.pptx"), recursive=True)
    files = pdfs + pptxs

    for f in files:
        path = Path(f)
        try:
            if path.suffix.lower() == ".pdf":
                items = _read_pdf(path)
            elif path.suffix.lower() == ".pptx":
                items = _read_pptx(path)
            else:
                continue
            for t, m in items:
                texts.append(t)
                metas.append(m)
        except Exception as e:
            print(f"[WARN] Skipping {path.name}: {e}")

    if not texts:
        raise RuntimeError("No chunks produced from dataset.")

    model = SentenceTransformer(EMBED_MODEL)
    embs = model.encode(texts, batch_size=64, show_progress_bar=True, normalize_embeddings=True)
    embs = np.asarray(embs, dtype="float32")
    dim = embs.shape[1]

    index = faiss.IndexFlatIP(dim)
    index.add(embs)
    faiss.write_index(index, str(INDEX_PATH))

    with open(CHUNKS_PATH, "w", encoding="utf-8") as f:
        for i, (t, m) in enumerate(zip(texts, metas)):
            rec = {"i": i, "text": t, "meta": m}
            f.write(json.dumps(rec, ensure_ascii=False) + "\n")

    with open(META_PATH, "w", encoding="utf-8") as f:
        json.dump({"dim": dim, "count": len(texts)}, f)

    return {"vectors": len(texts), "dim": dim, "files": len(files)}

# ---- search ----
_model = None
_index = None
_cache = None

def _load():
    global _model, _index, _cache
    if _model is None:
        _model = SentenceTransformer(EMBED_MODEL)
    if _index is None:
        _index = faiss.read_index(str(INDEX_PATH))
    if _cache is None:
        _cache = []
        with open(CHUNKS_PATH, "r", encoding="utf-8") as f:
            for line in f:
                _cache.append(json.loads(line))
    return _model, _index, _cache

def search(query: str, k: int = 5) -> List[Dict]:
    model, index, cache = _load()
    q = model.encode([query], normalize_embeddings=True)
    q = np.asarray(q, dtype="float32")
    D, I = index.search(q, k)
    results = []
    for score, idx in zip(D[0].tolist(), I[0].tolist()):
        rec = cache[idx]
        results.append({
            "score": float(score),
            "text": rec["text"],
            "meta": rec["meta"],
            "i": rec["i"],
        })
    return results
